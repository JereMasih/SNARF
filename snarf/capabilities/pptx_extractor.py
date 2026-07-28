import io

from pptx import Presentation

from snarf.capabilities.base import Capability


class PptxExtractor(Capability):
    name = "pptx_extractor"

    @property
    def available(self) -> bool:
        return True

    def extract_text(self, pptx_bytes: bytes) -> str:
        presentation = Presentation(io.BytesIO(pptx_bytes))
        parts = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text:
                    parts.append(shape.text_frame.text)
        return "\n".join(parts).strip()
