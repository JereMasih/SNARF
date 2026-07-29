import io

from fpdf import FPDF
from openpyxl import Workbook
from pptx import Presentation

from snarf.capabilities.base import Capability


class DocumentBuilder(Capability):
    """Genera bytes de documentos reales, todo local, sin credenciales ni
    costo — subirlos a Drive (y opcionalmente convertirlos a un formato
    nativo de Google) es responsabilidad de GoogleDrive.upload_file."""

    name = "document_builder"

    @property
    def available(self) -> bool:
        return True

    def build_markdown(self, title: str, content: str) -> bytes:
        text = f"# {title}\n\n{content}" if title else content
        return text.encode("utf-8")

    def build_pdf(self, title: str, content: str) -> bytes:
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 16)
        pdf.multi_cell(0, 10, title)
        pdf.ln(4)
        pdf.set_font("Helvetica", "", 11)
        pdf.multi_cell(0, 6, content)
        return bytes(pdf.output())

    def build_pptx(self, title: str, slides: list[dict]) -> bytes:
        presentation = Presentation()
        title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        title_slide.shapes.title.text = title
        body_layout = presentation.slide_layouts[1]
        for slide_data in slides:
            slide = presentation.slides.add_slide(body_layout)
            slide.shapes.title.text = slide_data.get("title", "")
            slide.placeholders[1].text_frame.text = slide_data.get("body", "")
        buf = io.BytesIO()
        presentation.save(buf)
        return buf.getvalue()

    def build_xlsx(self, title: str, rows: list[list]) -> bytes:
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = (title or "Hoja1")[:31]  # Excel limita el nombre de hoja a 31 caracteres
        for row in rows:
            sheet.append(row)
        buf = io.BytesIO()
        workbook.save(buf)
        return buf.getvalue()
