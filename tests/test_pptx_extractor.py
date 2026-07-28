import io

from pptx import Presentation

from snarf.capabilities.pptx_extractor import PptxExtractor


def make_pptx_bytes() -> bytes:
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = "título real de la diapositiva"
    slide.placeholders[1].text = "contenido real del cuerpo"
    buf = io.BytesIO()
    presentation.save(buf)
    return buf.getvalue()


def test_extract_text_includes_slide_title_and_body():
    text = PptxExtractor().extract_text(make_pptx_bytes())
    assert "título real de la diapositiva" in text
    assert "contenido real del cuerpo" in text


def test_is_always_available():
    assert PptxExtractor().available is True
